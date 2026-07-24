from pathlib import Path
from typing import List, Any
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader
from langchain_community.document_loaders import Docx2txtLoader
from langchain_core.documents import Document
from forum.services.rag_pipeline.video_loader import load_video_documents, VIDEO_EXTENSIONS
 
def load_all_documents(data_dir: str) -> List[Any]:
    data_path = Path(data_dir).resolve()
    print(f"[DEBUG] Data path: {data_path}")
    documents = []
    all_files = [f for f in data_path.glob('**/*') if f.is_file()]
 
    all_files = [f for f in all_files if not f.name.startswith('~$')]
 
    # Video files are handled separately by video_loader.py (transcription
    # + timestamp-aware chunking), so they're excluded from the generic loop below.
    non_video_files = [f for f in all_files if f.suffix.lower() not in VIDEO_EXTENSIONS]
    print(f"[DEBUG] found {len(non_video_files)} non-video files: {[f.name for f in non_video_files]}")
 
    for file_path in non_video_files:
        print(f"[DEBUG] loading file: {file_path}")
        try:
            if file_path.suffix == ".pdf":
                loader = PyPDFLoader(str(file_path))
                documents.extend(loader.load())
 
            elif file_path.suffix == ".txt":
                loader = TextLoader(str(file_path), encoding="utf-8")
                documents.extend(loader.load())
 
            elif file_path.suffix == ".docx":
                loader = Docx2txtLoader(str(file_path))
                documents.extend(loader.load())
 
            elif file_path.suffix == ".csv":
                loader = CSVLoader(str(file_path))
                documents.extend(loader.load())
 
            elif file_path.suffix == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path))
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text = f"Sheet: {sheet}\n"
                    for row in ws.iter_rows(values_only=True):
                        row_text = ' | '.join(str(c) for c in row if c is not None)
                        if row_text.strip():
                            text += row_text + "\n"
                    documents.append(Document(
                        page_content=text,
                        metadata={"source": str(file_path), "sheet": sheet}
                    ))
 
            elif file_path.suffix in [".pptx", ".ppt"]:
              if file_path.suffix == ".ppt":
                  try:
                      import comtypes.client
                      pptx_path = file_path.with_suffix(".pptx")
                      print(f"[INFO] Converting {file_path.name} to .pptx...")
                      powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                      powerpoint.Visible = 1
                      deck = powerpoint.Presentations.Open(str(file_path.resolve()))
                      deck.SaveAs(str(pptx_path.resolve()), 24)  # 24 = pptx format
                      deck.Close()
                      powerpoint.Quit()
                      file_path = pptx_path
                      print(f"[INFO] Converted to {pptx_path.name}")
                  except Exception as e:
                      print(f"[ERROR] Could not convert {file_path.name}: {e}")
                      continue
 
              from pptx import Presentation
              prs = Presentation(str(file_path))
              for i, slide in enumerate(prs.slides):
                  text = ""
                  for shape in slide.shapes:
                      if hasattr(shape, "text") and shape.text.strip():
                          text += shape.text.strip() + "\n"
                  if text.strip():
                      documents.append(Document(
                          page_content=text,
                          metadata={"source": str(file_path), "slide": i + 1}
                      ))
 
            elif file_path.suffix == ".json":
                import json
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                documents.append(Document(
                    page_content=str(data),
                    metadata={"source": str(file_path)}
                ))
 
            else:
                print(f"[SKIP] Unsupported file type: {file_path.suffix}")
                continue
 
            print(f"[DEBUG] Loaded from {file_path.name}")
 
        except Exception as e:
            print(f"[ERROR] Failed to load {file_path.name}: {e}")
 
    # Local video files: transcribe + chunk with timestamps, then merge in.
    video_docs = load_video_documents(str(data_path))
    documents.extend(video_docs)
 
    print(f"[INFO] Total documents loaded: {len(documents)}")
    return documents