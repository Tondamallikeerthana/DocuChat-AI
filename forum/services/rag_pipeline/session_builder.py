import os
from pathlib import Path
from typing import List

from langchain_core.documents import Document

from forum.services.rag_pipeline.vectorstore import FaissVectorStore

SESSION_STORE_ROOT = "session_stores"


def load_single_file(file_path: str) -> List[Document]:
    """Load one uploaded file into Document objects. Mirrors the
    per-extension logic in data_loader.load_all_documents, but for a
    single path instead of scanning a whole directory."""
    p = Path(file_path)
    suffix = p.suffix.lower()
    docs: List[Document] = []

    if suffix == ".pdf":
        from langchain_community.document_loaders import PyPDFLoader
        docs = PyPDFLoader(str(p)).load()

    elif suffix == ".txt":
        from langchain_community.document_loaders import TextLoader
        docs = TextLoader(str(p), encoding="utf-8").load()

    elif suffix == ".docx":
        from langchain_community.document_loaders import Docx2txtLoader
        docs = Docx2txtLoader(str(p)).load()

    elif suffix == ".csv":
        from langchain_community.document_loaders import CSVLoader
        docs = CSVLoader(str(p)).load()

    elif suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(p))
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
            if text.strip():
                docs.append(Document(page_content=text, metadata={"source": str(p), "slide": i + 1}))

    else:
        raise ValueError(f"Unsupported file type: {suffix}")

    if not docs:
        raise ValueError(f"No readable content found in {p.name}")

    return docs


def _session_dir(session_id: int) -> str:
    return os.path.join(SESSION_STORE_ROOT, f"session_{session_id}")


def build_session_store_from_documents(session_id: int, documents: List[Document]) -> FaissVectorStore:
    store = FaissVectorStore(persist_dir=_session_dir(session_id))
    store.build_from_documents(documents)
    return store


def load_session_store(session_id: int) -> FaissVectorStore:
    store = FaissVectorStore(persist_dir=_session_dir(session_id))
    store.load()
    return store